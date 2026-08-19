"""Slayt çizim ilkelleri / Slide rendering primitives.

Her yerleşim bloğu (`cover`, `quad`, `stats`, `panels`, `rows`, `flow`,
`closing`) tema nesnesinden renk alır; hiçbir renk burada sabit yazılmaz.
Böylece koyu ve açık varyantlar aynı koddan üretilir.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any, Iterable, Sequence

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

from presentation_theme import (
    ACCENT_CYCLE,
    BODY_Y,
    CARD_GAP_X,
    CARD_GAP_Y,
    CARD_W,
    CONTENT_W,
    FONT,
    MARGIN,
    SUB_Y,
    TILE_GAP,
    TILE_H,
    TILE_W,
    TITLE_Y,
    Theme,
)

ALIGN = {"left": PP_ALIGN.LEFT, "center": PP_ALIGN.CENTER, "right": PP_ALIGN.RIGHT}
ANCHOR = {"top": MSO_ANCHOR.TOP, "middle": MSO_ANCHOR.MIDDLE, "bottom": MSO_ANCHOR.BOTTOM}


def rgb(value: str) -> RGBColor:
    return RGBColor.from_string(value)


# ---------------------------------------------------------------------------
# Temel şekiller
# ---------------------------------------------------------------------------
def panel(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    fill: str,
    *,
    radius: float = 0.06,
    outline: str | None = None,
) -> Any:
    """Yuvarlatılmış dikdörtgen zemin."""
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    if outline:
        shape.line.color.rgb = rgb(outline)
        shape.line.width = Pt(0.75)
    else:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def circle(slide: Any, x: float, y: float, d: float, fill: str) -> Any:
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill)
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def text(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    runs: str | Sequence[tuple],
    *,
    size: float = 13,
    color: str = "FFFFFF",
    bold: bool = False,
    align: str = "left",
    anchor: str = "middle",
    spacing: float = 1.0,
    wrap: bool = True,
) -> Any:
    """Metin kutusu.

    `runs` düz metin ya da (metin, boyut, renk, kalın) demetlerinden oluşan bir
    dizi olabilir. Demet biçimi, emoji ile başlığı farklı puntoda yazmak için
    kullanılır (referans destede olduğu gibi).
    """
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = wrap
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = ANCHOR[anchor]

    items: Iterable[tuple]
    if isinstance(runs, str):
        items = [(runs, size, color, bold)]
    else:
        items = runs

    para = frame.paragraphs[0]
    para.alignment = ALIGN[align]
    if spacing != 1.0:
        para.line_spacing = spacing
    for chunk, csize, ccolor, cbold in items:
        run = para.add_run()
        run.text = chunk
        run.font.name = FONT
        run.font.size = Pt(csize)
        run.font.bold = cbold
        run.font.color.rgb = rgb(ccolor)
    return box


def paragraphs(
    slide: Any,
    x: float,
    y: float,
    w: float,
    h: float,
    lines: Sequence[str],
    *,
    size: float = 13,
    color: str = "FFFFFF",
    bold: bool = False,
    spacing: float = 1.16,
    space_after: float = 4,
) -> Any:
    """Çok satırlı metin kutusu — her öğe ayrı paragraf."""
    box = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    frame = box.text_frame
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = 0
    frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.TOP
    for index, line in enumerate(lines):
        para = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        para.line_spacing = spacing
        para.space_after = Pt(space_after)
        run = para.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = rgb(color)
    return box


def icon_badge(slide: Any, x: float, y: float, d: float, glyph: str, fill: str) -> None:
    """Renkli daire + ortasında emoji."""
    circle(slide, x, y, d, fill)
    text(
        slide,
        x,
        y,
        d,
        d,
        [(glyph, d * 32, "000000", False)],
        align="center",
        anchor="middle",
    )


# ---------------------------------------------------------------------------
# Slayt iskeleti
# ---------------------------------------------------------------------------
def new_slide(prs: Any, theme: Theme) -> Any:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = rgb(theme.bg)
    return slide


def header(slide: Any, theme: Theme, icon: str, title: str, subtitle: str | None) -> None:
    text(
        slide,
        MARGIN,
        TITLE_Y,
        CONTENT_W,
        0.75,
        [(f"{icon}  ", 26, theme.ink, False), (title, 29, theme.ink, True)],
        anchor="middle",
    )
    if subtitle:
        text(
            slide,
            MARGIN + 0.06,
            SUB_Y,
            CONTENT_W - 0.06,
            0.44,
            subtitle,
            size=13,
            color=theme.muted,
            anchor="top",
        )


BOTTOM = 7.12  # not şeridinin bittiği çizgi — altında 0.38" boşluk kalır


def note_height(message: str) -> float:
    """Not şeridi yüksekliğini metin uzunluğundan tahmin eder.

    Şeridi kalan boşluğu doldurmaya bırakmak, kısa notlarda devasa boş bir
    kutu üretiyordu; bu yüzden yükseklik metne göre sınırlandırılır.
    """
    length = len(message)
    if length > 200:
        return 1.10
    if length > 120:
        return 0.90
    return 0.74


def body_metrics(data: dict, lang: str) -> tuple[float, float | None]:
    """Gövde için kullanılabilir yüksekliği ve not şeridinin y'sini verir.

    Not şeridi her zaman slaytın altına sabitlenir; gövde blokları kalan
    alana göre küçülür. Bu, kartların not şeridiyle çakışmasını önler.
    """
    note = data.get("note")
    if not note:
        return 6.95 - BODY_Y, None
    y = BOTTOM - note_height(note[lang])
    return y - 0.18 - BODY_Y, y


def draw_note(slide: Any, theme: Theme, data: dict, lang: str, y: float | None) -> None:
    if y is None:
        return
    message = data["note"][lang]
    note_band(slide, theme, y, note_height(message), message, data.get("note_accent", "teal"))


def note_band(slide: Any, theme: Theme, y: float, h: float, message: str, accent: str) -> None:
    """Alt bilgi şeridi — sol kenarında renkli çubuk."""
    panel(slide, MARGIN, y, CONTENT_W, h, theme.card)
    bar = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(MARGIN + 0.16), Inches(y + 0.16), Inches(0.07), Inches(h - 0.32)
    )
    bar.adjustments[0] = 0.5
    bar.fill.solid()
    bar.fill.fore_color.rgb = rgb(theme.accent(accent))
    bar.line.fill.background()
    bar.shadow.inherit = False
    text(
        slide,
        MARGIN + 0.42,
        y + 0.10,
        CONTENT_W - 0.75,
        h - 0.20,
        message,
        size=13,
        color=theme.muted,
        anchor="middle",
        spacing=1.14,
    )


# ---------------------------------------------------------------------------
# Yerleşim blokları
# ---------------------------------------------------------------------------
def render_cover(slide: Any, theme: Theme, data: dict, lang: str) -> None:
    badge = data["badge"]
    panel(slide, 4.92, 1.02, 3.50, 1.15, theme.card, radius=0.14)
    text(
        slide,
        4.92,
        1.02,
        3.50,
        1.15,
        [(f"{badge['icon']}  ", 34, theme.ink, False), (badge["text"], 34, theme.teal, True)],
        align="center",
    )
    text(
        slide,
        1.00,
        2.48,
        11.33,
        1.05,
        data["title"][lang],
        size=44,
        color=theme.ink,
        bold=True,
        align="center",
    )
    text(slide, 1.00, 3.58, 11.33, 0.55, data["subtitle"][lang], size=19, color=theme.blue, align="center")
    text(slide, 1.00, 4.32, 11.33, 0.50, data["tagline"][lang], size=17, color=theme.gold, align="center")

    chips = data["chips"]
    chip_w, gap = 2.72, 0.16
    total = len(chips) * chip_w + (len(chips) - 1) * gap
    start = (13.333 - total) / 2
    for index, chip in enumerate(chips):
        cx = start + index * (chip_w + gap)
        panel(slide, cx, 5.30, chip_w, 0.58, theme.tile, radius=0.22)
        text(slide, cx, 5.30, chip_w, 0.58, chip[lang], size=12, color=theme.ink, align="center")

    text(slide, 1.00, 6.62, 11.33, 0.42, data["footer"][lang], size=11, color=theme.muted, align="center")


def render_quad(slide: Any, theme: Theme, data: dict, lang: str) -> None:
    """2x2 özellik kartı ızgarası."""
    items = data["items"]
    avail, note_y = body_metrics(data, lang)
    card_h = (avail - CARD_GAP_Y) / 2

    for index, item in enumerate(items):
        col, row = index % 2, index // 2
        x = MARGIN + col * (CARD_W + CARD_GAP_X)
        y = BODY_Y + row * (card_h + CARD_GAP_Y)
        accent = item.get("accent", ACCENT_CYCLE[index % len(ACCENT_CYCLE)])
        panel(slide, x, y, CARD_W, card_h, theme.card)
        icon_badge(slide, x + 0.22, y + 0.22, 0.52, item["icon"], theme.chip(accent))
        text(
            slide,
            x + 0.88,
            y + 0.22,
            CARD_W - 1.10,
            0.52,
            item["title"][lang],
            size=15,
            color=theme.ink,
            bold=True,
        )
        text(
            slide,
            x + 0.26,
            y + 0.90,
            CARD_W - 0.52,
            card_h - 1.10,
            item["body"][lang],
            size=12.5,
            color=theme.muted,
            anchor="top",
            spacing=1.16,
        )
    draw_note(slide, theme, data, lang, note_y)


def render_stats(slide: Any, theme: Theme, data: dict, lang: str) -> None:
    """4 sütunlu sayı kutuları (1 veya 2 satır) + isteğe bağlı not şeridi."""
    tiles = data["tiles"]
    rows = (len(tiles) + 3) // 4
    for index, tile in enumerate(tiles):
        col, row = index % 4, index // 4
        x = MARGIN + col * (TILE_W + TILE_GAP)
        y = 1.72 + row * (TILE_H + 0.30)
        accent = tile.get("accent", ACCENT_CYCLE[index % len(ACCENT_CYCLE)])
        panel(slide, x, y, TILE_W, TILE_H, theme.tile)
        text(
            slide,
            x,
            y + 0.12,
            TILE_W,
            0.80,
            tile["value"],
            size=31,
            color=theme.accent(accent),
            bold=True,
            align="center",
        )
        text(
            slide,
            x + 0.16,
            y + 0.92,
            TILE_W - 0.32,
            0.54,
            tile["label"][lang],
            size=10.5,
            color=theme.muted,
            align="center",
            anchor="top",
            spacing=1.08,
        )
    _ = rows
    draw_note(slide, theme, data, lang, body_metrics(data, lang)[1])


def render_panels(slide: Any, theme: Theme, data: dict, lang: str) -> None:
    """Yan yana iki panel; her panelde başlık + madde listesi."""
    left, right = data["left"], data["right"]
    ratio = data.get("ratio", 0.5)
    gap = 0.28
    width_l = (CONTENT_W - gap) * ratio
    width_r = CONTENT_W - gap - width_l
    height, note_y = body_metrics(data, lang)

    for spec, x, width in ((left, MARGIN, width_l), (right, MARGIN + width_l + gap, width_r)):
        accent = spec.get("accent", "teal")
        panel(slide, x, BODY_Y, width, height, theme.card)
        icon_badge(slide, x + 0.26, BODY_Y + 0.24, 0.50, spec["icon"], theme.chip(accent))
        text(
            slide,
            x + 0.90,
            BODY_Y + 0.24,
            width - 1.10,
            0.50,
            spec["title"][lang],
            size=15.5,
            color=theme.accent(accent),
            bold=True,
        )
        paragraphs(
            slide,
            x + 0.30,
            BODY_Y + 0.92,
            width - 0.60,
            height - 1.12,
            spec["lines"][lang],
            size=12.5,
            color=theme.muted,
            spacing=1.18,
            space_after=6,
        )
    draw_note(slide, theme, data, lang, note_y)


def render_rows(slide: Any, theme: Theme, data: dict, lang: str) -> None:
    """Tek geniş panel içinde iki sütuna dağılmış simge + başlık + açıklama satırları."""
    items = data["items"]
    height, note_y = body_metrics(data, lang)
    panel(slide, MARGIN, BODY_Y, CONTENT_W, height, theme.card)

    per_col = (len(items) + 1) // 2
    col_w = (CONTENT_W - 0.90) / 2
    # Satırlar panele yayılmak yerine sabit yükseklikte tutulup dikeyde
    # ortalanır; aksi hâlde kısa metinlerde panelin altı boş kalıyordu.
    row_h = min(1.26, (height - 0.50) / per_col)
    # Son satır kendi yuvasının tamamını doldurmaz; ortalarken yalnızca gerçek
    # içerik yüksekliği sayılır, aksi hâlde panelin altı üstünden geniş kalır.
    ink = (per_col - 1) * row_h + 0.95
    top = BODY_Y + max(0.28, (height - ink) / 2)

    for index, item in enumerate(items):
        col, row = index // per_col, index % per_col
        x = MARGIN + 0.30 + col * (col_w + 0.30)
        y = top + row * row_h
        accent = item.get("accent", ACCENT_CYCLE[index % len(ACCENT_CYCLE)])
        icon_badge(slide, x, y + 0.05, 0.42, item["icon"], theme.chip(accent))
        text(
            slide, x + 0.58, y, col_w - 0.62, 0.34, item["title"][lang], size=13.5, color=theme.ink, bold=True
        )
        text(
            slide,
            x + 0.58,
            y + 0.34,
            col_w - 0.62,
            row_h - 0.44,
            item["body"][lang],
            size=11.5,
            color=theme.muted,
            anchor="top",
            spacing=1.12,
        )
    draw_note(slide, theme, data, lang, note_y)


def render_flow(slide: Any, theme: Theme, data: dict, lang: str) -> None:
    """Numaralı yatay akış adımları."""
    steps = data["steps"]
    count = len(steps)
    gap = 0.20
    width = (CONTENT_W - (count - 1) * gap) / count
    height, note_y = body_metrics(data, lang)

    # İç içerik (rozet + başlık + gövde) sabit yüksekliktedir ve kart içinde
    # dikeyde ortalanır; kartın altında ölü boşluk kalmaz.
    # Gerçek içerik yüksekliği (rozet + başlık + 2-4 satır gövde) render
    # ölçümüyle ~2.15" çıkıyor; kartı bu değere göre ortalıyoruz.
    inner = 2.15
    pad = max(0.24, (height - inner) / 2)

    for index, step in enumerate(steps):
        x = MARGIN + index * (width + gap)
        accent = step.get("accent", ACCENT_CYCLE[index % len(ACCENT_CYCLE)])
        panel(slide, x, BODY_Y, width, height, theme.card)
        circle(slide, x + width / 2 - 0.24, BODY_Y + pad, 0.48, theme.chip(accent))
        text(
            slide,
            x + width / 2 - 0.24,
            BODY_Y + pad,
            0.48,
            0.48,
            str(index + 1),
            size=17,
            color=theme.ink,
            bold=True,
            align="center",
        )
        text(
            slide,
            x + 0.16,
            BODY_Y + pad + 0.62,
            width - 0.32,
            0.62,
            step["title"][lang],
            size=13,
            color=theme.accent(accent),
            bold=True,
            align="center",
            anchor="top",
            spacing=1.08,
        )
        text(
            slide,
            x + 0.18,
            BODY_Y + pad + 1.30,
            width - 0.36,
            inner - 1.30,
            step["body"][lang],
            size=11,
            color=theme.muted,
            align="center",
            anchor="top",
            spacing=1.14,
        )
        if index < count - 1:
            text(
                slide,
                x + width - 0.06,
                BODY_Y + height / 2 - 0.22,
                0.32,
                0.44,
                "›",
                size=22,
                color=theme.line,
                align="center",
            )
    draw_note(slide, theme, data, lang, note_y)


def render_closing(slide: Any, theme: Theme, data: dict, lang: str) -> None:
    text(
        slide,
        1.00,
        1.55,
        11.33,
        0.95,
        data["title"][lang],
        size=32,
        color=theme.ink,
        bold=True,
        align="center",
    )
    text(
        slide,
        1.60,
        2.55,
        10.13,
        0.60,
        data["subtitle"][lang],
        size=16,
        color=theme.blue,
        align="center",
        spacing=1.2,
    )

    items = data["items"]
    gap = 0.22
    width = (CONTENT_W - (len(items) - 1) * gap) / len(items)
    for index, item in enumerate(items):
        x = MARGIN + index * (width + gap)
        accent = item.get("accent", ACCENT_CYCLE[index % len(ACCENT_CYCLE)])
        panel(slide, x, 3.35, width, 1.70, theme.card)
        icon_badge(slide, x + width / 2 - 0.25, 3.55, 0.50, item["icon"], theme.chip(accent))
        text(
            slide,
            x + 0.14,
            4.18,
            width - 0.28,
            0.72,
            item["title"][lang],
            size=12.5,
            color=theme.ink,
            align="center",
            anchor="top",
            spacing=1.1,
        )
    note_band(slide, theme, 5.30, 0.92, data["note"][lang], data.get("note_accent", "gold"))
    text(slide, 1.00, 6.50, 11.33, 0.45, data["footer"][lang], size=11, color=theme.muted, align="center")


# ---------------------------------------------------------------------------
# Ekran görüntüsü blokları
# ---------------------------------------------------------------------------
SHOT_ROOT = Path(__file__).resolve().parent.parent / "sunum" / "ekranlar"
_CACHE = SHOT_ROOT / "_olcekli"

# Slaytta 9 inç genişliğe yerleşen bir görsel için 1600 piksel fazlasıyla
# yeterlidir; yakalama 2x ölçekle (3200 px) yapıldığı için gömmeden önce
# küçültülür. Aksi hâlde dört deste toplamda ~30 MB büyürdü.
_EMBED_WIDTH = 1600


def shot_path(theme: Theme, lang: str, name: str) -> Path | None:
    """Tema ve dile uygun ekran görüntüsünü döndürür, gerekirse küçülterek.

    Koyu deste koyu arayüzü, baskı destesi açık arayüzü gösterir; böylece
    slayt ile ekran görüntüsü aynı görsel dili konuşur.
    """
    mode = "dark" if theme.key == "ekran" else "light"
    source = SHOT_ROOT / f"{mode}_{lang}" / f"{name}.png"
    if not source.exists():
        return None

    target = _CACHE / f"{mode}_{lang}" / f"{name}.png"
    if target.exists() and target.stat().st_mtime >= source.stat().st_mtime:
        return target

    try:
        from PIL import Image

        target.parent.mkdir(parents=True, exist_ok=True)
        with Image.open(source) as image:
            if image.width > _EMBED_WIDTH:
                ratio = _EMBED_WIDTH / image.width
                image = image.resize((_EMBED_WIDTH, round(image.height * ratio)), Image.LANCZOS)
            image.save(target, "PNG", optimize=True)
        return target
    except Exception:
        # Küçültme başarısız olursa özgün dosya kullanılır — slayt yine üretilir.
        return source


def _place_shot(
    slide: Any,
    theme: Theme,
    path: Path,
    x: float,
    y: float,
    w: float,
    h: float,
) -> None:
    """Ekran görüntüsünü çerçeveli biçimde yerleştirir."""
    panel(slide, x - 0.06, y - 0.06, w + 0.12, h + 0.12, theme.line, radius=0.04)
    picture = slide.shapes.add_picture(str(path), Inches(x), Inches(y), Inches(w), Inches(h))
    picture.shadow.inherit = False


def render_shot(slide: Any, theme: Theme, data: dict, lang: str) -> None:
    """Tek büyük ekran görüntüsü + yanında açıklama maddeleri."""
    path = shot_path(theme, lang, data["shot"])
    shot_w, shot_h = 9.00, 5.63
    shot_y = 1.62

    if path is None:
        # Görsel yoksa slayt sessizce bozulmasın; durum açıkça yazılsın.
        panel(slide, MARGIN, shot_y, shot_w, shot_h, theme.card)
        text(
            slide,
            MARGIN,
            shot_y,
            shot_w,
            shot_h,
            "Ekran görüntüsü bulunamadı — tools/capture_screens.py çalıştırın",
            size=14,
            color=theme.muted,
            align="center",
        )
    else:
        _place_shot(slide, theme, path, MARGIN, shot_y, shot_w, shot_h)

    # Açıklama sütunu
    px, pw = MARGIN + shot_w + 0.25, CONTENT_W - shot_w - 0.25
    panel(slide, px, shot_y, pw, shot_h, theme.card)
    points = data.get("points", [])
    row_h = min(1.15, (shot_h - 0.50) / max(len(points), 1))
    top = shot_y + max(0.28, (shot_h - len(points) * row_h) / 2)

    for index, item in enumerate(points):
        accent = item.get("accent", ACCENT_CYCLE[index % len(ACCENT_CYCLE)])
        iy = top + index * row_h
        circle(slide, px + 0.24, iy + 0.04, 0.30, theme.accent(accent))
        text(
            slide,
            px + 0.66,
            iy - 0.02,
            pw - 0.90,
            0.34,
            item["title"][lang],
            size=11.5,
            color=theme.ink,
            bold=True,
        )
        text(
            slide,
            px + 0.66,
            iy + 0.30,
            pw - 0.90,
            row_h - 0.40,
            item["body"][lang],
            size=10,
            color=theme.muted,
            anchor="top",
            spacing=1.12,
        )


def render_gallery(slide: Any, theme: Theme, data: dict, lang: str) -> None:
    """Birden çok ekranı ızgara hâlinde gösterir."""
    items = data["items"]
    columns = data.get("columns", 3)
    rows = (len(items) + columns - 1) // columns

    gap = 0.22
    cell_w = (CONTENT_W - (columns - 1) * gap) / columns
    shot_h = cell_w / 1.6
    caption_h = 0.34
    cell_h = shot_h + caption_h

    total_h = rows * cell_h + (rows - 1) * gap
    top = BODY_Y + max(0.05, (6.95 - BODY_Y - total_h) / 2)

    for index, item in enumerate(items):
        col, row = index % columns, index // columns
        x = MARGIN + col * (cell_w + gap)
        y = top + row * (cell_h + gap)
        path = shot_path(theme, lang, item["shot"])
        if path is not None:
            _place_shot(slide, theme, path, x, y, cell_w, shot_h)
        else:
            panel(slide, x, y, cell_w, shot_h, theme.card)
        text(
            slide,
            x,
            y + shot_h + 0.04,
            cell_w,
            caption_h,
            item["title"][lang],
            size=11,
            color=theme.muted,
            align="center",
        )


RENDERERS = {
    "cover": render_cover,
    "quad": render_quad,
    "stats": render_stats,
    "panels": render_panels,
    "rows": render_rows,
    "flow": render_flow,
    "shot": render_shot,
    "gallery": render_gallery,
    "closing": render_closing,
}
